#CAN UPLOAD DOCUMENT FOR READING, CAN DO ADVANCED PROMPTING IN BOTH VOICE AND TEXT
import os
import google.generativeai as genai
import speech_recognition as sr
import pyttsx3
from tkinter import Tk, filedialog
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image, ImageEnhance, ImageOps
import pytesseract

# ONLY PRODUCTIVE AI NO ENTERTAINMENT ETC
# Fetch API key from the environment variable
#MUST BE CALLED IN ENVIORNMENT VARIABLE
genai.configure(api_key=os.environ["GENERATIVE_AI_API_KEY"])

# Generation configuration for the model
generation_config = {
    "temperature": 1.2,  # Increase for more variety in responses
    "top_p": 0.9,
    "top_k": 50,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}

# Initialize the GenerativeModel with your configuration
model = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=generation_config,
)

# Start the chat session without an initial history to avoid recalling past interactions
chat_session = model.start_chat(history=[])

# Initialize recognizer for speech recognition
recognizer = sr.Recognizer()

# Adjust recognizer settings
recognizer.energy_threshold = 4000  # Increase energy threshold for better sensitivity
recognizer.dynamic_energy_threshold = True  # Enable dynamic energy thresholding

# Initialize pyttsx3 for text-to-speech (optional, remove if not needed)
engine = pyttsx3.init()
voices = engine.getProperty('voices')

# Set voice to a female voice, if available
for voice in voices:
    if "female" in voice.name.lower() or "woman" in voice.name.lower():
        engine.setProperty('voice', voice.id)
        break

print("Welcome to the AI Assistant! Type 'exit' to quit.")

# Function to extract text from a PDF
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from a DOCX file
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from a PPTX file
def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from an image (with preprocessing)
def extract_text_from_image(image_path):
    image = Image.open(image_path)
    
    # Preprocess the image (convert to grayscale and enhance contrast)
    image = ImageOps.grayscale(image)  # Convert image to grayscale
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2)  # Increase contrast

    # Extract text using Tesseract
    text = pytesseract.image_to_string(image)
    return text

# Loop to continuously take queries from the user
while True:
    # Ask the user for input mode
    mode = input("Would you like to 'speak', 'type', or 'upload' a file? (type 'exit' to quit): ").strip().lower()

    if mode == "exit":
        print("Exiting the chat. Goodbye!")
        break

    elif mode == "speak":
        print("Listening... Please say your question or command.")
        try:
            # Use the microphone to capture speech
            with sr.Microphone() as source:
                print("Adjusting for ambient noise, please wait...")
                recognizer.adjust_for_ambient_noise(source, duration=1)  # Adjust for ambient noise
                audio_data = recognizer.listen(source, timeout=5)  # Increase listen timeout
                # Convert speech to text
                input_text = recognizer.recognize_google(audio_data)
                print(f"You said: {input_text}")
        except sr.UnknownValueError:
            print("Sorry, I didn't understand what you said. Please try again.")
            continue
        except sr.RequestError:
            print("Speech service is unavailable. Please try typing instead.")
            continue
        except sr.WaitTimeoutError:
            print("Listening timed out while waiting for phrase to start. Please try again.")
            continue

    elif mode == "type":
        # Get typed input
        input_text = input("Enter your question or command: ")

    elif mode == "upload":
        # Use tkinter file dialog to upload a file
        Tk().withdraw()  # Hide the root window
        file_path = filedialog.askopenfilename(
            title="Select a file",
            filetypes=[("PDF files", "*.pdf"), ("Word files", "*.docx"), ("PowerPoint files", "*.pptx"), ("Image files", "*.jpg;*.jpeg;*.png")]
        )

        if not file_path:
            print("No file selected.")
            continue

        # Extract text based on file type
        if file_path.endswith(".pdf"):
            input_text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".docx"):
            input_text = extract_text_from_docx(file_path)
        elif file_path.endswith(".pptx"):
            input_text = extract_text_from_pptx(file_path)
        elif file_path.endswith((".jpg", ".jpeg", ".png")):
            input_text = extract_text_from_image(file_path)
        else:
            print("Unsupported file type.")
            continue

        # Display extracted text preview
        print(f"Extracted text from the file: {input_text[:500]}...")  # Show the first 500 characters for review

    else:
        print("Invalid option. Please type 'speak', 'type', or 'upload'.")
        continue

    # Send the input text to the model and get a response
    response = chat_session.send_message(input_text)

    # Print the model's response
    response_text = response.text
    print("AI:", response_text)

    # Save conversation to history
    chat_session.history.append({"role": "user", "parts": [input_text]})
    chat_session.history.append({"role": "model", "parts": [response_text]})

    # Ask if the user wants to switch to another mode after processing the file
    next_action = input("\nWould you like to switch to 'speak', 'type', 'upload' again, or 'exit'? ").strip().lower()

    if next_action == "exit":
        print("Exiting the chat. Goodbye!")
        break
    elif next_action in ["speak", "type", "upload"]:
        continue
    else:
        print("Invalid input. Exiting the chat.")
        break
